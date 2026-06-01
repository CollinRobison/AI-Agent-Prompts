#!/usr/bin/env python3
"""Generate a themed one-pager slide from a template .pptx and a data file."""

from __future__ import annotations

import argparse
import csv
import json
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


# Conservative defaults used when a template does not expose theme details clearly.
DEFAULT_BG = "FFFFFF"
DEFAULT_TEXT = "1F1F1F"
DEFAULT_ACCENTS = ["2F5597", "4472C4", "70AD47", "FFC000", "ED7D31"]
DEFAULT_HEADING_FONT = "Calibri"
DEFAULT_BODY_FONT = "Calibri"
MAX_TABLE_COLUMNS = 4
MAX_TABLE_ROWS = 6
MIN_RECORDS_FOR_LINE_CHART = 3


@dataclass
class ThemeProfile:
    """Theme elements extracted from the template deck."""

    background_hex: str = DEFAULT_BG
    text_hex: str = DEFAULT_TEXT
    accent_hexes: list[str] = field(default_factory=lambda: DEFAULT_ACCENTS.copy())
    heading_font: str = DEFAULT_HEADING_FONT
    body_font: str = DEFAULT_BODY_FONT
    slide_width: int = 0
    slide_height: int = 0
    logo_positions: list[tuple[int, int, int, int]] = field(default_factory=list)


@dataclass
class NormalizedData:
    """Common data shape used by the slide renderer."""

    records: list[dict[str, Any]]
    summary: list[str] = field(default_factory=list)
    title: str | None = None


def _compute_layout(theme: ThemeProfile) -> dict[str, int]:
    """Build a proportional one-pager layout that fits any slide dimensions."""
    width = theme.slide_width
    height = theme.slide_height

    margin = int(width * 0.04)
    gutter = int(width * 0.03)
    title_top = int(height * 0.03)
    title_height = int(height * 0.1)

    left_width = int(width * 0.46)
    right_left = margin + left_width + gutter
    right_width = max(int(width * 0.2), width - right_left - margin)

    bullets_top = title_top + title_height + int(height * 0.015)
    bullets_height = int(height * 0.3)

    chart_top = bullets_top
    chart_height = int(height * 0.25)
    second_chart_top = chart_top + chart_height + int(height * 0.02)

    table_top = int(height * 0.67)
    table_height = max(int(height * 0.16), height - table_top - margin)

    return {
        "title_left": margin,
        "title_top": title_top,
        "title_width": width - (2 * margin),
        "title_height": title_height,
        "bullets_left": margin,
        "bullets_top": bullets_top,
        "bullets_width": left_width,
        "bullets_height": bullets_height,
        "chart_left": right_left,
        "chart_top": chart_top,
        "chart_width": right_width,
        "chart_height": chart_height,
        "chart2_top": second_chart_top,
        "table_left": margin,
        "table_top": table_top,
        "table_width": width - (2 * margin),
        "table_height": table_height,
    }


def _rgb_to_hex(rgb: RGBColor | None) -> str | None:
    """Convert python-pptx RGBColor to an uppercase 6-char hex string."""
    if rgb is None:
        return None
    value = str(rgb)
    if len(value) == 6:
        return value.upper()
    return None


def _extract_shape_color_hex(shape: Any) -> str | None:
    """Extract direct RGB fill color from a shape when present."""
    fill = getattr(shape, "fill", None)
    if not fill:
        return None
    try:
        fore_color = getattr(fill, "fore_color", None)
    except TypeError:
        # Some shapes are explicitly set to no-fill and raise on fore_color access.
        return None
    if not fore_color:
        return None
    return _rgb_to_hex(getattr(fore_color, "rgb", None))


def _extract_text_colors_and_fonts(shape: Any) -> tuple[list[str], list[str]]:
    """Extract text colors and fonts from runs in a text-bearing shape."""
    colors: list[str] = []
    fonts: list[str] = []
    text_frame = getattr(shape, "text_frame", None)
    if not text_frame:
        return colors, fonts

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font is None:
                continue
            color_hex = _rgb_to_hex(getattr(getattr(run.font, "color", None), "rgb", None))
            if color_hex:
                colors.append(color_hex)
            if run.font.name:
                fonts.append(run.font.name)
    return colors, fonts


def extract_theme(template_path: Path) -> ThemeProfile:
    """Phase 1: inspect template deck and pull a reusable theme profile."""
    prs = Presentation(str(template_path))

    accent_candidates: list[str] = []
    text_candidates: list[str] = []
    font_candidates: list[str] = []
    logo_positions: list[tuple[int, int, int, int]] = []

    background_hex = None

    for slide in prs.slides:
        if background_hex is None:
            bg_fill = slide.background.fill
            try:
                bg_color = _rgb_to_hex(getattr(getattr(bg_fill, "fore_color", None), "rgb", None))
            except TypeError:
                bg_color = None
            if bg_color:
                background_hex = bg_color

        for shape in slide.shapes:
            # Capture potential logo placement to optionally preserve brand placement.
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                logo_positions.append((shape.left, shape.top, shape.width, shape.height))

            shape_color = _extract_shape_color_hex(shape)
            if shape_color:
                accent_candidates.append(shape_color)

            run_colors, run_fonts = _extract_text_colors_and_fonts(shape)
            text_candidates.extend(run_colors)
            font_candidates.extend(run_fonts)

    # De-duplicate while preserving order so first colors stay the strongest candidates.
    def unique(values: Iterable[str]) -> list[str]:
        seen: set[str] = set()
        output: list[str] = []
        for val in values:
            if val not in seen:
                seen.add(val)
                output.append(val)
        return output

    accent_unique = unique(accent_candidates)
    text_unique = unique(text_candidates)
    fonts_unique = unique(font_candidates)

    heading_font = fonts_unique[0] if fonts_unique else DEFAULT_HEADING_FONT
    body_font = fonts_unique[1] if len(fonts_unique) > 1 else heading_font

    return ThemeProfile(
        background_hex=background_hex or DEFAULT_BG,
        text_hex=text_unique[0] if text_unique else DEFAULT_TEXT,
        accent_hexes=accent_unique[:5] if accent_unique else DEFAULT_ACCENTS.copy(),
        heading_font=heading_font,
        body_font=body_font,
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        logo_positions=logo_positions[:3],
    )


def _read_json(path: Path) -> NormalizedData:
    with path.open("r", encoding="utf-8") as f:
        payload = json.load(f)

    if isinstance(payload, list):
        records = [item for item in payload if isinstance(item, dict)]
        return NormalizedData(records=records)

    if isinstance(payload, dict):
        records = payload.get("records") or payload.get("data")
        summary = payload.get("summary")
        title = payload.get("title")

        normalized_records: list[dict[str, Any]] = []
        if isinstance(records, list):
            normalized_records = [item for item in records if isinstance(item, dict)]
        elif all(not isinstance(v, (dict, list)) for v in payload.values()):
            # Flat object: represent as one-row table.
            normalized_records = [payload]

        normalized_summary = [str(x) for x in summary] if isinstance(summary, list) else []
        return NormalizedData(records=normalized_records, summary=normalized_summary, title=title)

    raise ValueError("Unsupported JSON structure. Expected a list or object.")


def _read_csv(path: Path) -> NormalizedData:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f)
        return NormalizedData(records=[dict(row) for row in reader])


def _read_excel(path: Path) -> NormalizedData:
    try:
        df = pd.read_excel(path)
    except Exception as exc:
        raise ValueError(f"Failed to read Excel file '{path.name}': {exc}") from exc
    records = df.where(pd.notna(df), None).to_dict(orient="records")
    return NormalizedData(records=records)


def _parse_markdown_table(lines: list[str]) -> list[dict[str, Any]]:
    """Parse the first markdown table found and convert it to records."""
    # Keep only table-like lines containing pipes.
    table_lines = [line.strip() for line in lines if "|" in line]
    if len(table_lines) < 2:
        return []

    header = [h.strip() for h in table_lines[0].strip("|").split("|")]
    records: list[dict[str, Any]] = []

    for line in table_lines[2:]:
        cells = [c.strip() for c in line.strip("|").split("|")]
        if len(cells) != len(header):
            continue
        row = {header[i]: cells[i] for i in range(len(header))}
        records.append(row)

    return records


def _read_markdown(path: Path) -> NormalizedData:
    lines = path.read_text(encoding="utf-8").splitlines()
    records = _parse_markdown_table(lines)
    summary = [line.lstrip("- ").strip() for line in lines if line.strip().startswith(("- ", "* "))]
    return NormalizedData(records=records, summary=summary)


def normalize_data(path: Path) -> NormalizedData:
    """Phase 2: parse supported formats and normalize them into list-of-dicts records."""
    suffix = path.suffix.lower()
    if suffix == ".json":
        data = _read_json(path)
    elif suffix == ".csv":
        data = _read_csv(path)
    elif suffix in {".xlsx", ".xls"}:
        data = _read_excel(path)
    elif suffix == ".md":
        data = _read_markdown(path)
    else:
        raise ValueError(f"Unsupported data format: {suffix}")

    if not data.records:
        raise ValueError("No records were parsed from the data source.")
    return data


def _to_rgb(hex_color: str) -> RGBColor:
    hex_color = (hex_color or DEFAULT_TEXT).lstrip("#")
    if len(hex_color) != 6:
        hex_color = DEFAULT_TEXT
    return RGBColor.from_string(hex_color.upper())


def _to_float(value: Any) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).replace(",", "").replace("%", "").strip()
    if not text:
        return None
    try:
        return float(text)
    except ValueError:
        return None


def _pick_label_and_value_fields(records: list[dict[str, Any]]) -> tuple[str, str] | None:
    if not records:
        return None

    keys = list(records[0].keys())
    if not keys:
        return None

    label_key = "label" if "label" in keys else keys[0]

    for key in keys:
        if key == label_key:
            continue
        if any(_to_float(row.get(key)) is not None for row in records):
            return label_key, key

    return None


def _build_summary(data: NormalizedData) -> list[str]:
    """Use provided summary when available, otherwise derive concise bullet points."""
    if data.summary:
        return data.summary[:5]

    records = data.records
    fields = _pick_label_and_value_fields(records)
    if not fields:
        return ["Data loaded successfully."]

    label_key, value_key = fields
    values = [(str(row.get(label_key, "")), _to_float(row.get(value_key))) for row in records]
    values = [(label, val) for label, val in values if val is not None]
    if not values:
        return ["Data loaded successfully."]

    max_label, max_value = max(values, key=lambda item: item[1])
    min_label, min_value = min(values, key=lambda item: item[1])
    avg_value = sum(v for _, v in values) / len(values)

    return [
        f"Highest {value_key}: {max_label} ({max_value:.2f})",
        f"Lowest {value_key}: {min_label} ({min_value:.2f})",
        f"Average {value_key}: {avg_value:.2f}",
    ]


def _create_chart_image(
    records: list[dict[str, Any]],
    label_key: str,
    value_key: str,
    chart_type: str,
    accent_hexes: list[str],
    output_path: Path,
) -> None:
    labels = [str(row.get(label_key, "")) for row in records]
    values = [_to_float(row.get(value_key)) for row in records]

    filtered = [(label, value) for label, value in zip(labels, values) if value is not None]
    if not filtered:
        raise ValueError("No numeric values available for chart generation.")

    labels, values = zip(*filtered)
    colors = [f"#{accent_hexes[i % len(accent_hexes)]}" for i in range(len(values))]

    plt.figure(figsize=(5.4, 2.8))
    if chart_type == "line":
        plt.plot(labels, values, marker="o", color=colors[0], linewidth=2)
    else:
        plt.bar(labels, values, color=colors)

    plt.xticks(rotation=25, ha="right", fontsize=8)
    plt.yticks(fontsize=8)
    plt.tight_layout()
    plt.savefig(output_path, dpi=160, bbox_inches="tight", transparent=True)
    plt.close()


def _set_slide_background(slide: Any, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _to_rgb(color_hex)


def _set_text_shape_style(shape: Any, font_name: str, color_hex: str, size_pt: int, bold: bool = False) -> None:
    has_text_frame = getattr(shape, "has_text_frame", None)
    if has_text_frame is False:
        return
    frame = getattr(shape, "text_frame", None)
    if frame is None:
        return
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = _to_rgb(color_hex)


def _add_title(slide: Any, text: str, theme: ThemeProfile, layout: dict[str, int]) -> None:
    title_shape = slide.shapes.add_textbox(
        layout["title_left"],
        layout["title_top"],
        layout["title_width"],
        layout["title_height"],
    )
    tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = theme.heading_font
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = _to_rgb(theme.text_hex)


def _add_bullets(slide: Any, bullets: list[str], theme: ThemeProfile, layout: dict[str, int]) -> None:
    bullet_shape = slide.shapes.add_textbox(
        layout["bullets_left"],
        layout["bullets_top"],
        layout["bullets_width"],
        layout["bullets_height"],
    )
    tf = bullet_shape.text_frame
    tf.clear()

    for idx, bullet in enumerate(bullets[:5]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = theme.body_font
        p.font.size = Pt(16)
        p.font.color.rgb = _to_rgb(theme.text_hex)


def _add_table(slide: Any, records: list[dict[str, Any]], theme: ThemeProfile, layout: dict[str, int]) -> None:
    if not records:
        return

    columns = list(records[0].keys())[:MAX_TABLE_COLUMNS]
    if not columns:
        return

    row_count = min(len(records), MAX_TABLE_ROWS)
    table_shape = slide.shapes.add_table(
        row_count + 1,
        len(columns),
        layout["table_left"],
        layout["table_top"],
        layout["table_width"],
        layout["table_height"],
    )
    table = table_shape.table

    for col_idx, col_name in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _to_rgb(theme.accent_hexes[0])
        _set_text_shape_style(cell, theme.body_font, DEFAULT_BG, 12, bold=True)

    for row_idx, row in enumerate(records[:row_count], start=1):
        for col_idx, col_name in enumerate(columns):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(row.get(col_name, ""))
            _set_text_shape_style(cell, theme.body_font, theme.text_hex, 11)


def generate_onepager(template_path: Path, data_path: Path, output_path: Path, title_override: str | None) -> None:
    """Phase 3: build a single-slide output deck styled from extracted theme + data."""
    theme = extract_theme(template_path)
    data = normalize_data(data_path)

    title = title_override or data.title or "Data Analysis One-Pager"
    bullets = _build_summary(data)

    prs = Presentation()
    prs.slide_width = theme.slide_width or prs.slide_width
    prs.slide_height = theme.slide_height or prs.slide_height
    theme.slide_width = prs.slide_width
    theme.slide_height = prs.slide_height
    layout = _compute_layout(theme)

    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    _set_slide_background(slide, theme.background_hex)
    _add_title(slide, title, theme, layout)
    _add_bullets(slide, bullets, theme, layout)

    fields = _pick_label_and_value_fields(data.records)
    if fields:
        label_key, value_key = fields
        with tempfile.TemporaryDirectory(prefix="onepager_charts_") as tmpdir:
            tmpdir_path = Path(tmpdir)
            bar_chart_path = tmpdir_path / "bar_chart.png"

            _create_chart_image(data.records, label_key, value_key, "bar", theme.accent_hexes, bar_chart_path)

            # Embed one or two charts in the right-hand side visual panel.
            slide.shapes.add_picture(
                str(bar_chart_path),
                layout["chart_left"],
                layout["chart_top"],
                layout["chart_width"],
                layout["chart_height"],
            )
            if len(data.records) >= MIN_RECORDS_FOR_LINE_CHART:
                line_chart_path = tmpdir_path / "line_chart.png"
                _create_chart_image(data.records, label_key, value_key, "line", theme.accent_hexes, line_chart_path)
                slide.shapes.add_picture(
                    str(line_chart_path),
                    layout["chart_left"],
                    layout["chart2_top"],
                    layout["chart_width"],
                    layout["chart_height"],
                )

    _add_table(slide, data.records, theme, layout)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Generate a themed one-pager slide from a template and data file.")
    parser.add_argument("--template", required=True, help="Path to source .pptx template")
    parser.add_argument("--data", required=True, help="Path to input data file (.json/.csv/.xlsx/.xls/.md)")
    parser.add_argument("--output", default="onepager.pptx", help="Path for generated one-pager .pptx")
    parser.add_argument("--title", default=None, help="Optional title override for the one-pager")
    return parser


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    template_path = Path(args.template)
    data_path = Path(args.data)
    output_path = Path(args.output)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    if template_path.suffix.lower() != ".pptx":
        raise ValueError("Template must be a .pptx file")
    if not data_path.exists():
        raise FileNotFoundError(f"Data file not found: {data_path}")

    generate_onepager(template_path, data_path, output_path, args.title)
    print(f"Generated one-pager: {output_path}")


if __name__ == "__main__":
    main()
