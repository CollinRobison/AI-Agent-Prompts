#!/usr/bin/env python3
"""Generate themed presentation slides from a template .pptx and a data file.

The script extracts the complete visual theme from the template's slide-master
XML (official color scheme slots dk1/lt1/accent1-6, major/minor fonts, slide
dimensions) then builds one or more slides that faithfully reproduce the brand.
Each slide type targets a different analytical purpose:

  Slide 1 – Executive Summary   title, KPI callouts, insight bullets, compact chart
  Slide 2 – Visual Analysis      large chart(s) with annotations and trend insights
  Slide 3 – Data Detail          full data table with secondary chart

Usage:
  python generate_onepager.py --template deck.pptx --data data.json \\
      --output report.pptx --title "Q1 Analysis" --slides auto
"""

from __future__ import annotations

import argparse
import csv
import json
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt


# ── Fallback defaults (used when template provides no usable signals) ──────────
DEFAULT_BG = "FFFFFF"
DEFAULT_TEXT = "1F1F1F"
DEFAULT_ACCENTS = ["2F5597", "4472C4", "70AD47", "FFC000", "ED7D31", "C00000"]
DEFAULT_HEADING_FONT = "Calibri"
DEFAULT_BODY_FONT = "Calibri"

MAX_TABLE_COLUMNS = 5
MAX_TABLE_ROWS_SUMMARY = 6
MAX_TABLE_ROWS_DETAIL = 15
MIN_RECORDS_FOR_LINE = 4

# Minimum records to justify an additional analysis/detail slide in auto mode.
AUTO_SLIDE_2_THRESHOLD = 8
AUTO_SLIDE_3_THRESHOLD = 18


# ── Data structures ────────────────────────────────────────────────────────────

@dataclass
class ThemeProfile:
    """Full visual theme extracted from the template's slide-master XML."""

    background_hex: str = DEFAULT_BG
    text_hex: str = DEFAULT_TEXT
    accent_hexes: list[str] = field(default_factory=lambda: DEFAULT_ACCENTS.copy())
    muted_bg_hex: str = "F2F2F2"    # alternating table rows, section dividers
    heading_font: str = DEFAULT_HEADING_FONT
    body_font: str = DEFAULT_BODY_FONT
    slide_width: int = 0
    slide_height: int = 0
    logo_positions: list[tuple[int, int, int, int]] = field(default_factory=list)
    is_dark_bg: bool = False


@dataclass
class NormalizedData:
    """Common data shape consumed by every slide renderer."""

    records: list[dict[str, Any]]
    summary: list[str] = field(default_factory=list)
    title: str | None = None


# ── Color helpers ──────────────────────────────────────────────────────────────

def _is_dark(hex_color: str) -> bool:
    """Return True when perceived luminance of hex_color is below 0.45."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255 < 0.45


def _lighten(hex_color: str, factor: float = 0.82) -> str:
    """Blend hex_color toward white by factor (0 = unchanged, 1 = white)."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return hex_color
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, int(r + (255 - r) * factor))
    g = min(255, int(g + (255 - g) * factor))
    b = min(255, int(b + (255 - b) * factor))
    return f"{r:02X}{g:02X}{b:02X}"


def _rgb_to_hex(rgb: RGBColor | None) -> str | None:
    if rgb is None:
        return None
    v = str(rgb)
    return v.upper() if len(v) == 6 else None


def _to_rgb(hex_color: str) -> RGBColor:
    h = (hex_color or DEFAULT_TEXT).lstrip("#")
    if len(h) != 6:
        h = DEFAULT_TEXT
    return RGBColor.from_string(h.upper())


# ── Theme extraction ───────────────────────────────────────────────────────────

def _read_xml_theme(prs: Presentation) -> tuple[dict[str, str], str, str]:
    """Extract color scheme and font scheme from slide-master theme XML.

    Returns (color_map, heading_font, body_font).
    color_map keys: dk1 lt1 dk2 lt2 accent1 … accent6
    Falls back to empty dict / DEFAULT_* strings on any error.
    """
    color_map: dict[str, str] = {}
    heading_font = DEFAULT_HEADING_FONT
    body_font = DEFAULT_BODY_FONT

    try:
        master = prs.slide_master
        for rel in master.part.rels.values():
            if not rel.reltype.endswith("/theme"):
                continue
            theme_elem = rel.target_part._element

            # Color scheme: the authoritative slot names PowerPoint respects.
            clr_scheme = theme_elem.find(".//" + qn("a:clrScheme"))
            if clr_scheme is not None:
                roles = ["dk1", "lt1", "dk2", "lt2",
                         "accent1", "accent2", "accent3",
                         "accent4", "accent5", "accent6"]
                for role in roles:
                    node = clr_scheme.find(qn(f"a:{role}"))
                    if node is None:
                        continue
                    srgb = node.find(qn("a:srgbClr"))
                    if srgb is not None:
                        val = srgb.get("val", "").upper()
                        if val:
                            color_map[role] = val
                        continue
                    sys_clr = node.find(qn("a:sysClr"))
                    if sys_clr is not None:
                        val = sys_clr.get("lastClr", "").upper()
                        if val:
                            color_map[role] = val

            # Font scheme: major = headings, minor = body.
            font_scheme = theme_elem.find(".//" + qn("a:fontScheme"))
            if font_scheme is not None:
                def _face(tag: str) -> str | None:
                    parent = font_scheme.find(qn(tag))
                    if parent is None:
                        return None
                    latin = parent.find(qn("a:latin"))
                    if latin is None:
                        return None
                    face = latin.get("typeface", "")
                    # Relative references start with '+'; skip them.
                    return face if (face and not face.startswith("+")) else None

                heading_font = _face("a:majorFont") or DEFAULT_HEADING_FONT
                body_font = _face("a:minorFont") or DEFAULT_BODY_FONT
            break  # only the first theme relationship is needed
    except Exception:
        pass  # degrade gracefully; defaults will be used

    return color_map, heading_font, body_font


def _scan_slide_colors(prs: Presentation) -> tuple[list[str], list[str], str | None]:
    """Scan slide shapes for fill/text colors as supplementary fallback signals.

    Returns (accent_candidates, text_candidates, background_hex_or_None).
    """
    accent_candidates: list[str] = []
    text_candidates: list[str] = []
    background_hex: str | None = None

    def _unique(vals: Iterable[str]) -> list[str]:
        seen: set[str] = set()
        out: list[str] = []
        for v in vals:
            if v not in seen:
                seen.add(v)
                out.append(v)
        return out

    for slide in prs.slides:
        if background_hex is None:
            try:
                bg_hex = _rgb_to_hex(
                    getattr(getattr(slide.background.fill, "fore_color", None), "rgb", None)
                )
                if bg_hex:
                    background_hex = bg_hex
            except TypeError:
                pass

        for shape in slide.shapes:
            fill = getattr(shape, "fill", None)
            if fill:
                try:
                    fc = getattr(fill, "fore_color", None)
                    if fc:
                        h = _rgb_to_hex(getattr(fc, "rgb", None))
                        if h:
                            accent_candidates.append(h)
                except TypeError:
                    pass

            tf = getattr(shape, "text_frame", None)
            if tf:
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.font:
                            h = _rgb_to_hex(
                                getattr(getattr(run.font, "color", None), "rgb", None)
                            )
                            if h:
                                text_candidates.append(h)

    return _unique(accent_candidates), _unique(text_candidates), background_hex


def extract_theme(template_path: Path) -> ThemeProfile:
    """Build a ThemeProfile from template, preferring slide-master XML signals."""
    prs = Presentation(str(template_path))

    color_map, heading_font, body_font = _read_xml_theme(prs)
    accent_scan, text_scan, bg_scan = _scan_slide_colors(prs)

    # ── Background (lt1 = official light/background slot) ────────────────────
    bg_hex = color_map.get("lt1") or bg_scan or DEFAULT_BG

    # ── Text (dk1 = official dark/text slot) ─────────────────────────────────
    text_hex = color_map.get("dk1") or (text_scan[0] if text_scan else DEFAULT_TEXT)

    # ── Accent palette (accent1–6 from XML, fill in from scan or defaults) ───
    xml_accents = [
        color_map[f"accent{i}"] for i in range(1, 7) if color_map.get(f"accent{i}")
    ]
    if not xml_accents:
        xml_accents = accent_scan[:6] if accent_scan else DEFAULT_ACCENTS.copy()
    while len(xml_accents) < 5:
        xml_accents.append(xml_accents[0])

    # ── Logo positions ────────────────────────────────────────────────────────
    logo_positions: list[tuple[int, int, int, int]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                logo_positions.append((shape.left, shape.top, shape.width, shape.height))
        if logo_positions:
            break

    is_dark = _is_dark(bg_hex)
    if is_dark and text_hex == DEFAULT_TEXT:
        text_hex = "FFFFFF"

    profile = ThemeProfile(
        background_hex=bg_hex,
        text_hex=text_hex,
        accent_hexes=xml_accents,
        muted_bg_hex=_lighten(xml_accents[0]) if xml_accents else "F2F2F2",
        heading_font=heading_font,
        body_font=body_font,
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        logo_positions=logo_positions[:3],
        is_dark_bg=is_dark,
    )

    print(
        f"[theme]  bg={profile.background_hex}  text={profile.text_hex}  "
        f"accents={profile.accent_hexes[:3]}  "
        f"fonts={profile.heading_font}/{profile.body_font}  dark_bg={profile.is_dark_bg}"
    )
    return profile


# ── Data ingestion ─────────────────────────────────────────────────────────────

def _read_json(path: Path) -> NormalizedData:
    with path.open("r", encoding="utf-8") as f:
        payload = json.load(f)

    if isinstance(payload, list):
        return NormalizedData(records=[r for r in payload if isinstance(r, dict)])

    if isinstance(payload, dict):
        records = payload.get("records") or payload.get("data")
        summary = payload.get("summary")
        title = payload.get("title")

        normalized: list[dict[str, Any]] = []
        if isinstance(records, list):
            normalized = [r for r in records if isinstance(r, dict)]
        elif all(not isinstance(v, (dict, list)) for v in payload.values()):
            normalized = [payload]

        return NormalizedData(
            records=normalized,
            summary=[str(x) for x in summary] if isinstance(summary, list) else [],
            title=title,
        )

    raise ValueError("Unsupported JSON structure. Expected a list or object.")


def _read_csv(path: Path) -> NormalizedData:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return NormalizedData(records=[dict(row) for row in csv.DictReader(f)])


def _read_excel(path: Path) -> NormalizedData:
    try:
        df = pd.read_excel(path)
    except Exception as exc:
        raise ValueError(f"Failed to read Excel file '{path.name}': {exc}") from exc
    return NormalizedData(records=df.where(pd.notna(df), None).to_dict(orient="records"))


def _read_markdown(path: Path) -> NormalizedData:
    lines = path.read_text(encoding="utf-8").splitlines()
    table_lines = [l.strip() for l in lines if "|" in l]
    records: list[dict[str, Any]] = []
    if len(table_lines) >= 2:
        headers = [h.strip() for h in table_lines[0].strip("|").split("|")]
        for line in table_lines[2:]:
            cells = [c.strip() for c in line.strip("|").split("|")]
            if len(cells) == len(headers):
                records.append({headers[i]: cells[i] for i in range(len(headers))})
    summary = [l.lstrip("- ").strip() for l in lines if l.strip().startswith(("- ", "* "))]
    return NormalizedData(records=records, summary=summary)


def normalize_data(path: Path) -> NormalizedData:
    """Parse supported formats and normalize into list-of-dict records."""
    suffix = path.suffix.lower()
    dispatch = {
        ".json": _read_json,
        ".csv": _read_csv,
        ".xlsx": _read_excel,
        ".xls": _read_excel,
        ".md": _read_markdown,
    }
    if suffix not in dispatch:
        raise ValueError(f"Unsupported data format: {suffix!r}")
    data = dispatch[suffix](path)
    if not data.records:
        raise ValueError("No records were parsed from the data source.")
    print(f"[data]   {len(data.records)} records  cols={list(data.records[0].keys())[:6]}")
    return data


# ── Numeric & analytical helpers ───────────────────────────────────────────────

def _to_float(value: Any) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).replace(",", "").replace("%", "").strip()
    try:
        return float(text) if text else None
    except ValueError:
        return None


def _pick_label_and_value_fields(records: list[dict[str, Any]]) -> tuple[str, str] | None:
    """Return (label_col, first_numeric_col) or None."""
    if not records:
        return None
    keys = list(records[0].keys())
    label_key = "label" if "label" in keys else keys[0]
    for key in keys:
        if key == label_key:
            continue
        if any(_to_float(row.get(key)) is not None for row in records):
            return label_key, key
    return None


def _all_numeric_columns(records: list[dict[str, Any]]) -> list[str]:
    """Return all column names with at least one numeric value."""
    if not records:
        return []
    return [
        k for k in records[0]
        if any(_to_float(row.get(k)) is not None for row in records)
    ]


def _format_number(value: float) -> str:
    """Human-friendly number for KPI display."""
    if abs(value) >= 1_000_000:
        return f"{value / 1_000_000:.1f}M"
    if abs(value) >= 1_000:
        return f"{value / 1_000:.1f}K"
    if value == int(value):
        return f"{int(value):,}"
    return f"{value:,.2f}"


def _extract_kpis(data: NormalizedData) -> list[tuple[str, str]]:
    """Return up to 3 (display_value, label) KPI tuples."""
    fields = _pick_label_and_value_fields(data.records)
    if not fields:
        return []
    label_key, value_key = fields
    pairs = [
        (str(r.get(label_key, "")), _to_float(r.get(value_key)))
        for r in data.records
    ]
    pairs = [(l, v) for l, v in pairs if v is not None]
    if not pairs:
        return []
    max_l, max_v = max(pairs, key=lambda x: x[1])
    _min_l, _min_v = min(pairs, key=lambda x: x[1])
    total = sum(v for _, v in pairs)
    avg = total / len(pairs)
    return [
        (_format_number(total), f"Total {value_key.replace('_', ' ').title()}"),
        (_format_number(max_v), f"Peak  ·  {max_l}"),
        (_format_number(avg), f"Average {value_key.replace('_', ' ').title()}"),
    ]


def _build_insight_bullets(data: NormalizedData) -> list[str]:
    """Derive analytical insight bullets. Uses user-supplied summary when present."""
    if data.summary:
        return data.summary[:6]

    fields = _pick_label_and_value_fields(data.records)
    if not fields:
        return ["▪ Data loaded successfully."]

    label_key, value_key = fields
    pairs = [
        (str(r.get(label_key, "")), _to_float(r.get(value_key)))
        for r in data.records
    ]
    pairs = [(l, v) for l, v in pairs if v is not None]
    if not pairs:
        return ["▪ Data loaded successfully."]

    max_l, max_v = max(pairs, key=lambda x: x[1])
    min_l, min_v = min(pairs, key=lambda x: x[1])
    avg = sum(v for _, v in pairs) / len(pairs)
    span_pct = ((max_v - min_v) / min_v * 100) if min_v != 0 else 0

    bullets = [
        f"▲  Peak {value_key.replace('_', ' ')}: {max_l} — {_format_number(max_v)}",
        f"▼  Lowest: {min_l} — {_format_number(min_v)}",
        f"◆  Average: {_format_number(avg)}",
    ]
    if span_pct > 5:
        bullets.append(f"↕  Range spans {span_pct:.1f}% from low to high")

    num_cols = _all_numeric_columns(data.records)
    if len(num_cols) > 1:
        bullets.append(
            f"📊  {len(num_cols)} metrics available: "
            + ", ".join(c.replace("_", " ") for c in num_cols[:3])
        )

    return bullets


def _infer_chart_type(records: list[dict[str, Any]], label_key: str) -> str:
    """Pick bar, line, or pie based on data characteristics."""
    n = len(records)
    time_kws = {
        "jan", "feb", "mar", "apr", "may", "jun", "jul", "aug",
        "sep", "oct", "nov", "dec", "q1", "q2", "q3", "q4",
        "week", "day", "month", "year", "quarter",
        "2020", "2021", "2022", "2023", "2024", "2025", "2026",
    }
    labels = [str(r.get(label_key, "")).lower() for r in records]
    is_temporal = any(
        any(kw in label for kw in time_kws) for label in labels
    )
    if n >= MIN_RECORDS_FOR_LINE and is_temporal:
        return "line"
    if n <= 7:
        return "pie"
    return "bar"


def _decide_slide_count(data: NormalizedData) -> int:
    """Auto-determine slide count from data size and column richness."""
    n = len(data.records)
    num_cols = len(_all_numeric_columns(data.records))
    if n >= AUTO_SLIDE_3_THRESHOLD or num_cols >= 4:
        return 3
    if n >= AUTO_SLIDE_2_THRESHOLD or num_cols >= 2:
        return 2
    return 1


# ── Chart rendering ────────────────────────────────────────────────────────────

def _create_chart_image(
    records: list[dict[str, Any]],
    label_key: str,
    value_key: str,
    chart_type: str,
    theme: ThemeProfile,
    chart_title: str,
    output_path: Path,
    figsize: tuple[float, float] = (6.0, 3.4),
) -> None:
    """Render a styled chart image matched to the extracted theme."""
    labels = [str(r.get(label_key, "")) for r in records]
    raw = [_to_float(r.get(value_key)) for r in records]
    pairs = [(l, v) for l, v in zip(labels, raw) if v is not None]
    if not pairs:
        raise ValueError(f"No numeric values in column '{value_key}' for chart.")
    labels_f, values_f = zip(*pairs)

    bg = f"#{theme.background_hex}"
    txt = f"#{theme.text_hex}"
    accents = [f"#{h}" for h in theme.accent_hexes]
    primary = accents[0] if accents else "#4472C4"
    grid_color = "#CCCCCC" if not theme.is_dark_bg else "#444444"
    spine_color = grid_color

    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)

    if chart_type == "pie":
        pie_colors = [accents[i % len(accents)] for i in range(len(values_f))]
        wedges, texts, autotexts = ax.pie(
            values_f,
            labels=labels_f,
            colors=pie_colors,
            autopct="%1.1f%%",
            startangle=90,
            pctdistance=0.75,
            textprops={"fontsize": 8, "color": txt},
        )
        for at in autotexts:
            at.set_color("white")
            at.set_fontweight("bold")
            at.set_fontsize(7)
        ax.axis("equal")

    elif chart_type == "line":
        x_idx = list(range(len(labels_f)))
        ax.plot(x_idx, values_f, marker="o", color=primary, linewidth=2.5,
                markersize=6, markerfacecolor="white", markeredgewidth=2,
                markeredgecolor=primary, zorder=3)
        ax.fill_between(x_idx, values_f, alpha=0.12, color=primary)
        for i, val in enumerate(values_f):
            fmt = _format_number(val)
            ax.annotate(fmt, xy=(i, val), xytext=(0, 9), textcoords="offset points",
                        ha="center", fontsize=7, color=txt, fontweight="bold")
        ax.set_xticks(x_idx)
        ax.set_xticklabels(labels_f, rotation=30, ha="right", fontsize=8, color=txt)
        ax.tick_params(axis="y", colors=txt, labelsize=8)
        ax.grid(axis="y", color=grid_color, linestyle="--", linewidth=0.5, zorder=0)
        ax.set_axisbelow(True)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color(spine_color)
        ax.spines["bottom"].set_color(spine_color)
        ax.set_ylabel(value_key.replace("_", " ").title(), fontsize=8, color=txt)

    else:  # bar
        bar_colors = [accents[i % len(accents)] for i in range(len(values_f))]
        x_idx = list(range(len(labels_f)))
        bars = ax.bar(x_idx, values_f, color=bar_colors, edgecolor="none", width=0.6, zorder=2)
        for bar, val in zip(bars, values_f):
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                bar.get_height(),
                _format_number(val),
                ha="center", va="bottom", fontsize=7, color=txt, fontweight="bold",
            )
        ax.set_xticks(x_idx)
        ax.set_xticklabels(labels_f, rotation=30, ha="right", fontsize=8, color=txt)
        ax.tick_params(axis="y", colors=txt, labelsize=8)
        ax.grid(axis="y", color=grid_color, linestyle="--", linewidth=0.5, zorder=0)
        ax.set_axisbelow(True)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color(spine_color)
        ax.spines["bottom"].set_color(spine_color)
        ax.set_ylabel(value_key.replace("_", " ").title(), fontsize=8, color=txt)

    ax.set_title(chart_title, fontsize=10, fontweight="bold", color=txt, pad=8)

    plt.tight_layout(pad=0.8)
    plt.savefig(output_path, dpi=160, bbox_inches="tight",
                facecolor=bg, edgecolor="none")
    plt.close()


# ── Slide-building primitives ──────────────────────────────────────────────────

def _new_blank_slide(prs: Presentation, theme: ThemeProfile) -> Any:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _to_rgb(theme.background_hex)
    return slide


def _add_title_block(
    slide: Any, text: str, theme: ThemeProfile,
    left: int, top: int, width: int, height: int,
    font_size: int = 28,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = theme.heading_font
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = _to_rgb(theme.text_hex)


def _add_accent_bar(
    slide: Any, theme: ThemeProfile,
    left: int, top: int, width: int, height: int = 0,
) -> None:
    """Thin colored horizontal rule under the title — brand accent stripe."""
    if height == 0:
        height = int(theme.slide_height * 0.006)
    bar = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = _to_rgb(theme.accent_hexes[0])
    bar.line.fill.background()


def _add_kpi_row(
    slide: Any, kpis: list[tuple[str, str]], theme: ThemeProfile,
    left: int, top: int, total_width: int, height: int,
) -> None:
    """Row of colored KPI callout boxes (value + label)."""
    if not kpis:
        return
    n = len(kpis)
    gutter = int(total_width * 0.02)
    box_w = (total_width - gutter * (n - 1)) // n
    on_dark_bg = _is_dark(theme.accent_hexes[0])
    txt_on_accent = "FFFFFF" if on_dark_bg else theme.text_hex

    for i, (value_str, label_str) in enumerate(kpis):
        bx = left + i * (box_w + gutter)
        accent = theme.accent_hexes[i % len(theme.accent_hexes)]

        box = slide.shapes.add_shape(1, bx, top, box_w, height)
        box.fill.solid()
        box.fill.fore_color.rgb = _to_rgb(accent)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        tf.margin_top = Pt(8)
        tf.margin_bottom = Pt(4)

        p_val = tf.paragraphs[0]
        p_val.alignment = PP_ALIGN.CENTER
        r_val = p_val.add_run()
        r_val.text = value_str
        r_val.font.name = theme.heading_font
        r_val.font.size = Pt(24)
        r_val.font.bold = True
        r_val.font.color.rgb = _to_rgb(txt_on_accent)

        p_lbl = tf.add_paragraph()
        p_lbl.alignment = PP_ALIGN.CENTER
        r_lbl = p_lbl.add_run()
        r_lbl.text = label_str
        r_lbl.font.name = theme.body_font
        r_lbl.font.size = Pt(9)
        r_lbl.font.color.rgb = _to_rgb(txt_on_accent)


def _add_bullets_block(
    slide: Any, bullets: list[str], theme: ThemeProfile,
    left: int, top: int, width: int, height: int,
    font_size: int = 14,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)

    for idx, bullet in enumerate(bullets[:7]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.space_before = Pt(4)
        p.font.name = theme.body_font
        p.font.size = Pt(font_size)
        p.font.color.rgb = _to_rgb(theme.text_hex)


def _add_table_block(
    slide: Any, records: list[dict[str, Any]], theme: ThemeProfile,
    left: int, top: int, width: int, height: int,
    max_rows: int = MAX_TABLE_ROWS_SUMMARY,
) -> None:
    if not records:
        return
    columns = list(records[0].keys())[:MAX_TABLE_COLUMNS]
    if not columns:
        return
    row_count = min(len(records), max_rows)

    tbl = slide.shapes.add_table(
        row_count + 1, len(columns), left, top, width, height
    ).table

    # Header row styled with primary accent.
    for ci, col in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = col.replace("_", " ").title()
        cell.fill.solid()
        cell.fill.fore_color.rgb = _to_rgb(theme.accent_hexes[0])
        tf = cell.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.name = theme.body_font
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = _to_rgb("FFFFFF")

    # Data rows with alternating muted background.
    for ri, row in enumerate(records[:row_count], start=1):
        row_bg = theme.muted_bg_hex if ri % 2 == 0 else theme.background_hex
        for ci, col in enumerate(columns):
            cell = tbl.cell(ri, ci)
            cell.text = str(row.get(col, "") or "")
            cell.fill.solid()
            cell.fill.fore_color.rgb = _to_rgb(row_bg)
            tf = cell.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.name = theme.body_font
                    run.font.size = Pt(10)
                    run.font.color.rgb = _to_rgb(theme.text_hex)


# ── Slide generators ───────────────────────────────────────────────────────────

def _add_summary_slide(
    prs: Presentation, theme: ThemeProfile,
    data: NormalizedData, title: str, tmpdir: Path,
) -> None:
    """Slide 1 — Executive Summary: KPI row, insight bullets, compact chart, mini-table."""
    slide = _new_blank_slide(prs, theme)
    W, H = theme.slide_width, theme.slide_height
    margin = int(W * 0.04)
    gutter = int(W * 0.025)

    # Title
    title_top = int(H * 0.03)
    title_h = int(H * 0.09)
    _add_title_block(slide, title, theme, margin, title_top, W - 2 * margin, title_h, font_size=26)

    # Accent bar
    bar_top = title_top + title_h + int(H * 0.005)
    _add_accent_bar(slide, theme, margin, bar_top, W - 2 * margin)

    # KPI row
    kpi_top = bar_top + int(H * 0.015)
    kpi_h = int(H * 0.14)
    kpis = _extract_kpis(data)
    if kpis:
        _add_kpi_row(slide, kpis, theme, margin, kpi_top, W - 2 * margin, kpi_h)

    # Two-column zone below KPIs: bullets left, chart right
    zone_top = kpi_top + kpi_h + int(H * 0.02)
    zone_h = int(H * 0.34)
    left_w = int(W * 0.44)
    right_left = margin + left_w + gutter
    right_w = W - right_left - margin

    bullets = _build_insight_bullets(data)
    _add_bullets_block(slide, bullets, theme, margin, zone_top, left_w, zone_h, font_size=13)

    fields = _pick_label_and_value_fields(data.records)
    if fields:
        label_key, value_key = fields
        chart_type = _infer_chart_type(data.records, label_key)
        chart_path = tmpdir / "summary_chart.png"
        _create_chart_image(
            data.records, label_key, value_key, chart_type, theme,
            value_key.replace("_", " ").title(), chart_path, figsize=(5.2, 3.0),
        )
        slide.shapes.add_picture(str(chart_path), right_left, zone_top, right_w, zone_h)

    # Compact data table at bottom
    tbl_top = zone_top + zone_h + int(H * 0.015)
    tbl_h = max(int(H * 0.17), H - tbl_top - margin)
    _add_table_block(slide, data.records, theme, margin, tbl_top, W - 2 * margin, tbl_h,
                     max_rows=MAX_TABLE_ROWS_SUMMARY)


def _add_analysis_slide(
    prs: Presentation, theme: ThemeProfile,
    data: NormalizedData, title: str, tmpdir: Path,
) -> None:
    """Slide 2 — Visual Analysis: large primary chart, secondary chart, trend annotation."""
    slide = _new_blank_slide(prs, theme)
    W, H = theme.slide_width, theme.slide_height
    margin = int(W * 0.04)
    gutter = int(W * 0.025)

    title_top = int(H * 0.03)
    title_h = int(H * 0.09)
    _add_title_block(slide, f"{title}  ·  Visual Analysis", theme,
                     margin, title_top, W - 2 * margin, title_h, font_size=22)
    bar_top = title_top + title_h + int(H * 0.005)
    _add_accent_bar(slide, theme, margin, bar_top, W - 2 * margin)

    fields = _pick_label_and_value_fields(data.records)
    if not fields:
        _add_bullets_block(slide, _build_insight_bullets(data), theme,
                           margin, bar_top + int(H * 0.04), W - 2 * margin,
                           int(H * 0.7), font_size=15)
        return

    label_key, value_key = fields
    primary_type = _infer_chart_type(data.records, label_key)
    secondary_type = "bar" if primary_type == "line" else "line"

    chart_top = bar_top + int(H * 0.04)
    left_chart_w = int(W * 0.55)
    right_chart_left = margin + left_chart_w + gutter
    right_chart_w = W - right_chart_left - margin
    chart_h = int(H * 0.45)

    primary_path = tmpdir / "analysis_primary.png"
    _create_chart_image(
        data.records, label_key, value_key, primary_type, theme,
        f"{value_key.replace('_', ' ').title()} by {label_key.replace('_', ' ').title()}",
        primary_path, figsize=(6.5, 3.8),
    )
    slide.shapes.add_picture(str(primary_path), margin, chart_top, left_chart_w, chart_h)

    # Secondary chart — only if data has enough records for a second view.
    if len(data.records) >= MIN_RECORDS_FOR_LINE:
        secondary_path = tmpdir / "analysis_secondary.png"
        _create_chart_image(
            data.records, label_key, value_key, secondary_type, theme,
            f"Trend — {value_key.replace('_', ' ').title()}",
            secondary_path, figsize=(4.5, 3.8),
        )
        slide.shapes.add_picture(str(secondary_path), right_chart_left, chart_top,
                                 right_chart_w, chart_h)

    # Insight annotation strip below charts.
    anno_top = chart_top + chart_h + int(H * 0.02)
    anno_h = H - anno_top - margin
    bullets = _build_insight_bullets(data)
    _add_bullets_block(slide, bullets[:4], theme, margin, anno_top, W - 2 * margin,
                       anno_h, font_size=13)


def _add_detail_slide(
    prs: Presentation, theme: ThemeProfile,
    data: NormalizedData, title: str, tmpdir: Path,
) -> None:
    """Slide 3 — Data Detail: full table with a secondary chart alongside."""
    slide = _new_blank_slide(prs, theme)
    W, H = theme.slide_width, theme.slide_height
    margin = int(W * 0.04)
    gutter = int(W * 0.025)

    title_top = int(H * 0.03)
    title_h = int(H * 0.09)
    _add_title_block(slide, f"{title}  ·  Data Detail", theme,
                     margin, title_top, W - 2 * margin, title_h, font_size=22)
    bar_top = title_top + title_h + int(H * 0.005)
    _add_accent_bar(slide, theme, margin, bar_top, W - 2 * margin)

    tbl_top = bar_top + int(H * 0.025)
    tbl_w = int(W * 0.55)
    chart_left = margin + tbl_w + gutter
    chart_w = W - chart_left - margin
    content_h = H - tbl_top - margin

    _add_table_block(slide, data.records, theme, margin, tbl_top, tbl_w, content_h,
                     max_rows=MAX_TABLE_ROWS_DETAIL)

    fields = _pick_label_and_value_fields(data.records)
    if fields:
        label_key, value_key = fields
        num_cols = _all_numeric_columns(data.records)
        # If multiple numeric columns exist, chart the second one for variety.
        second_value = num_cols[1] if len(num_cols) > 1 else value_key
        detail_chart_path = tmpdir / "detail_chart.png"
        _create_chart_image(
            data.records, label_key, second_value, "bar", theme,
            second_value.replace("_", " ").title(),
            detail_chart_path, figsize=(4.5, 4.2),
        )
        slide.shapes.add_picture(str(detail_chart_path), chart_left, tbl_top, chart_w, content_h)


# ── Main generation pipeline ───────────────────────────────────────────────────

def generate_presentation(
    template_path: Path,
    data_path: Path,
    output_path: Path,
    title_override: str | None,
    slide_count: int | str,
) -> None:
    """Build a fully themed presentation and save to output_path."""
    theme = extract_theme(template_path)
    data = normalize_data(data_path)

    title = title_override or data.title or "Data Analysis"

    # Resolve slide count.
    if slide_count == "auto":
        n_slides = _decide_slide_count(data)
    else:
        n_slides = int(slide_count)
    n_slides = max(1, min(n_slides, 3))
    print(f"[slides] generating {n_slides} slide(s) — title: {title!r}")

    prs = Presentation()
    prs.slide_width = theme.slide_width or prs.slide_width
    prs.slide_height = theme.slide_height or prs.slide_height
    theme.slide_width = prs.slide_width
    theme.slide_height = prs.slide_height

    with tempfile.TemporaryDirectory(prefix="onepager_charts_") as tmpdir:
        tmpdir_path = Path(tmpdir)

        _add_summary_slide(prs, theme, data, title, tmpdir_path)
        if n_slides >= 2:
            _add_analysis_slide(prs, theme, data, title, tmpdir_path)
        if n_slides >= 3:
            _add_detail_slide(prs, theme, data, title, tmpdir_path)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Generate themed presentation slides from a .pptx template and data file."
    )
    parser.add_argument("--template", required=True, help="Path to source .pptx template")
    parser.add_argument("--data", required=True,
                        help="Path to input data file (.json/.csv/.xlsx/.xls/.md)")
    parser.add_argument("--output", default="onepager.pptx",
                        help="Output .pptx path (default: onepager.pptx)")
    parser.add_argument("--title", default=None, help="Optional title override")
    parser.add_argument(
        "--slides", default="auto",
        help="Number of slides to generate: 1, 2, 3, or auto (default: auto)."
             " auto picks 1 for small datasets, 2 for medium, 3 for large.",
    )
    return parser


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    template_path = Path(args.template)
    data_path = Path(args.data)
    output_path = Path(args.output)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    if template_path.suffix.lower() != ".pptx":
        raise ValueError("Template must be a .pptx file.")
    if not data_path.exists():
        raise FileNotFoundError(f"Data file not found: {data_path}")
    if args.slides not in {"auto", "1", "2", "3"}:
        raise ValueError("--slides must be 1, 2, 3, or auto.")

    generate_presentation(template_path, data_path, output_path, args.title, args.slides)
    print(f"[done]   saved → {output_path}")


if __name__ == "__main__":
    main()
