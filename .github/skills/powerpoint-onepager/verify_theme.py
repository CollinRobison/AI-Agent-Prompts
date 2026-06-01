#!/usr/bin/env python3
"""Verify that a generated .pptx matches the visual theme of its source template.

Usage:
    python verify_theme.py --template deck.pptx --output report.pptx

Exit code 0 = all checks passed, 1 = one or more checks failed.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn


def _parse_theme_blob(prs: Presentation) -> tuple[etree._Element | None, str]:
    """Return (clrScheme_element, color_scheme_name) by parsing the theme blob.

    Reads the raw bytes of the theme XML part (works even when python-pptx
    returns a plain Part without _element for the theme relationship).
    Returns (None, '') if no theme is found.
    """
    try:
        master = prs.slide_master
        for rel in master.part.rels.values():
            if rel.reltype.endswith("/theme"):
                root = etree.fromstring(rel.target_part.blob)
                clr = root.find(".//" + qn("a:clrScheme"))
                name = clr.get("name", "") if clr is not None else ""
                return clr, name
    except Exception:
        pass
    return None, ""


def _extract_colors(clr_scheme: etree._Element | None) -> dict[str, str]:
    """Return slot→hex dict from a clrScheme element."""
    if clr_scheme is None:
        return {}
    color_map: dict[str, str] = {}
    for role in ["dk1", "lt1", "dk2", "lt2",
                 "accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]:
        node = clr_scheme.find(qn(f"a:{role}"))
        if node is None:
            continue
        srgb = node.find(qn("a:srgbClr"))
        if srgb is not None:
            val = srgb.get("val", "").upper()
            if val:
                color_map[role] = val
            continue
        sys_clr = node.find(qn("a:sysClr"))
        if sys_clr is not None:
            val = sys_clr.get("lastClr", "").upper()
            if val:
                color_map[role] = val
    return color_map


def _extract_fonts(prs: Presentation) -> tuple[str, str]:
    """Return (heading_font, body_font) by parsing the theme blob."""
    try:
        master = prs.slide_master
        for rel in master.part.rels.values():
            if rel.reltype.endswith("/theme"):
                root = etree.fromstring(rel.target_part.blob)
                fs = root.find(".//" + qn("a:fontScheme"))
                if fs is None:
                    return "", ""

                def _face(tag: str) -> str:
                    parent = fs.find(qn(tag))
                    if parent is None:
                        return ""
                    latin = parent.find(qn("a:latin"))
                    if latin is None:
                        return ""
                    face = latin.get("typeface", "")
                    return face if (face and not face.startswith("+")) else ""

                return _face("a:majorFont"), _face("a:minorFont")
    except Exception:
        pass
    return "", ""


def verify(template_path: Path, output_path: Path) -> bool:
    tmpl = Presentation(str(template_path))
    out = Presentation(str(output_path))

    all_pass = True

    def check(label: str, passed: bool, detail: str = "") -> None:
        nonlocal all_pass
        status = "PASS" if passed else "FAIL"
        msg = f"[verify] {label:<22} {status}"
        if detail:
            msg += f"  {detail}"
        print(msg)
        if not passed:
            all_pass = False

    # ── Slide dimensions ──────────────────────────────────────────────────────
    w_ok = tmpl.slide_width == out.slide_width
    h_ok = tmpl.slide_height == out.slide_height
    check("slide dimensions", w_ok and h_ok,
          f"{out.slide_width} x {out.slide_height}")

    # ── Theme color scheme name + serialized XML fingerprint ──────────────────
    t_clr, t_name = _parse_theme_blob(tmpl)
    o_clr, o_name = _parse_theme_blob(out)

    if t_clr is not None and o_clr is not None:
        t_xml = etree.tostring(t_clr, encoding="unicode")
        o_xml = etree.tostring(o_clr, encoding="unicode")
        master_ok = t_xml == o_xml
        detail = f"scheme='{o_name}'" if master_ok else f"template='{t_name}'  output='{o_name}'"
    else:
        master_ok = False
        detail = "could not read theme XML from one or both files"
    check("slide master theme", master_ok, detail)

    # ── Accent colors ─────────────────────────────────────────────────────────
    t_colors = _extract_colors(t_clr)
    o_colors = _extract_colors(o_clr)
    t_acc = [t_colors.get(f"accent{i}", "") for i in range(1, 7) if t_colors.get(f"accent{i}")]
    o_acc = [o_colors.get(f"accent{i}", "") for i in range(1, 7) if o_colors.get(f"accent{i}")]
    check("accent colors", t_acc == o_acc,
          f"template={t_acc[:3]}  output={o_acc[:3]}")

    # ── Background (lt1) ──────────────────────────────────────────────────────
    t_bg = t_colors.get("lt1", "")
    o_bg = o_colors.get("lt1", "")
    check("background (lt1)", t_bg == o_bg, f"bg={o_bg or '(none)'}")

    # ── Fonts ─────────────────────────────────────────────────────────────────
    t_h, t_b = _extract_fonts(tmpl)
    o_h, o_b = _extract_fonts(out)
    check("fonts", t_h == o_h and t_b == o_b,
          f"heading={o_h or '(none)'}  body={o_b or '(none)'}")

    # ── Slide count ───────────────────────────────────────────────────────────
    n = len(out.slides)
    check("slide count", n >= 1, f"{n} slide(s)")

    verdict = "PASS — output matches template theme" if all_pass else "FAIL — see checks above"
    print(f"[result] {verdict}")
    return all_pass


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Verify a generated .pptx matches its source template's theme."
    )
    parser.add_argument("--template", required=True, help="Path to source .pptx template")
    parser.add_argument("--output", required=True, help="Path to generated .pptx to verify")
    args = parser.parse_args()

    template_path = Path(args.template)
    output_path = Path(args.output)

    if not template_path.exists():
        print(f"[error] Template not found: {template_path}", file=sys.stderr)
        sys.exit(1)
    if not output_path.exists():
        print(f"[error] Output file not found: {output_path}", file=sys.stderr)
        sys.exit(1)

    passed = verify(template_path, output_path)
    sys.exit(0 if passed else 1)


if __name__ == "__main__":
    main()
